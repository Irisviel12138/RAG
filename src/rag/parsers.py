"""Document parsers for PDF, DOCX, and PPTX files."""

from __future__ import annotations

import re
import unicodedata
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import List


@dataclass
class ParsedDocument:
    doc_id: str
    content: str
    source_type: str


def _normalize_text(text: str) -> str:
    # Unicode normalization helps keep math symbols and latin letters stable.
    text = unicodedata.normalize("NFKC", text)

    # Remove common replacement chars produced by PDF extraction.
    text = text.replace("�", "")

    # Normalize common OCR/PDF bullets/dashes.
    text = text.replace("•", "- ").replace("·", "- ")

    # Heuristic fixes for broken derivatives in teaching slides/PDFs.
    text = text.replace("d f", "df").replace("d x", "dx")

    # Collapse extra spaces but preserve line breaks for structure.
    lines = [re.sub(r"[ \t]+", " ", ln).strip() for ln in text.splitlines()]
    lines = [ln for ln in lines if ln]

    return "\n".join(lines)


def _read_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(data))
    pages: List[str] = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(f"[Page {i}]\n{_normalize_text(text)}")
    return "\n\n".join(pages)


def _read_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(BytesIO(data))
    parts: List[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(_normalize_text(text))
    return "\n".join(parts)


def _read_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    slides: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        chunk: List[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            text = text.strip() if text else ""
            if text:
                chunk.append(_normalize_text(text))
        if chunk:
            slides.append(f"[Slide {idx}]\n" + "\n".join(chunk))
    return "\n\n".join(slides)


def parse_document(filename: str, data: bytes) -> ParsedDocument:
    ext = Path(filename).suffix.lower()
    doc_id = Path(filename).stem

    if ext == ".pdf":
        content = _read_pdf(data)
        return ParsedDocument(doc_id=doc_id, content=content, source_type="pdf")

    if ext == ".docx":
        content = _read_docx(data)
        return ParsedDocument(doc_id=doc_id, content=content, source_type="docx")

    if ext == ".pptx":
        content = _read_pptx(data)
        return ParsedDocument(doc_id=doc_id, content=content, source_type="pptx")

    raise ValueError(f"Unsupported file type: {ext}. Use .pdf, .docx, or .pptx")
